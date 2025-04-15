from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import pandas as pd
import os

def insert_table(slide, df, left, top, width, height):
    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table

    # Set column names
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = str(col_name)

    # Set row values
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

    return table

def insert_image(slide, image_path, left, top, width=None, height=None):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        print(f"Image not found: {image_path}")

def add_heading(slide, text, left, top, font_size=24):
    textbox = slide.shapes.add_textbox(left, top, Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.bold = True
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0, 0, 0)

def generate_presentation(template_path, output_path, table_entries, plot_entries):
    prs = Presentation(template_path)

    for entry in table_entries:
        slide = prs.slides[entry['slide_number']]
        add_heading(slide, entry['heading'], Inches(0.5), Inches(0.2))
        insert_table(slide, entry['dataframe'], Inches(0.5), Inches(1), Inches(9), Inches(4))

    for entry in plot_entries:
        slide = prs.slides[entry['slide_number']]
        add_heading(slide, entry['heading'], Inches(0.5), Inches(0.2))
        insert_image(slide, entry['image_path'], Inches(0.5), Inches(1), Inches(6), Inches(4))

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
